"""Shared slide builder for the two ISHS decks.

Geometry and colours are lifted from the existing house deck
(core/frontend/docs/SmartGate-Monitoring-Ritase.pptx) so these sit alongside it
rather than looking like a different product: 13.33x7.5in, #07090D ground, amber
eyebrow rule, Segoe UI titles, Consolas for anything machine-ish.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from PIL import Image

BG = RGBColor(0x07, 0x09, 0x0D)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
TEXT = RGBColor(0xF8, 0xFA, 0xFC)
DIM = RGBColor(0x6B, 0x76, 0x88)
SECOND = RGBColor(0xA3, 0xAE, 0xC2)
FRAME = RGBColor(0x2A, 0x31, 0x42)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)

SANS, MONO = "Segoe UI", "Consolas"


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    return prs


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _text(slide, left, top, width, height, text, *, size, color, font=SANS,
          bold=False, align=PP_ALIGN.LEFT, spacing=1.0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _rect(slide, left, top, width, height, fill, *, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = 0.06
        except (IndexError, KeyError):
            pass
    return sh


def title_slide(prs, eyebrow, title, subtitle, tag):
    slide = _blank(prs)
    _rect(slide, 0.55, 2.30, 0.07, 1.55, AMBER, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, 0.80, 2.28, 10.0, 0.3, eyebrow, size=10.5, color=AMBER, font=MONO)
    _text(slide, 0.76, 2.62, 11.6, 0.9, title, size=40, color=TEXT, bold=True)
    _text(slide, 0.78, 3.62, 10.4, 0.9, subtitle, size=15, color=SECOND, spacing=1.25)
    _text(slide, 0.78, 5.10, 8.0, 0.3, tag, size=10, color=DIM, font=MONO)
    return slide


def _head(slide, eyebrow, title):
    _rect(slide, 0.55, 0.52, 0.06, 0.62, AMBER, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, 0.78, 0.50, 9.0, 0.24, eyebrow, size=9.5, color=AMBER, font=MONO)
    _text(slide, 0.75, 0.72, 11.5, 0.55, title, size=25, color=TEXT, bold=True)


def _foot(slide, footer, page, total):
    _text(slide, 0.55, 6.95, 6.5, 0.25, footer, size=9, color=DIM, font=MONO)
    _text(slide, 11.20, 6.95, 1.60, 0.25, f"{page} / {total}", size=9, color=DIM,
          font=MONO, align=PP_ALIGN.RIGHT)


def _bullets(slide, left, top, width, height, notes, *, size=12.5, color=AMBER):
    body = "\n".join(f"▪  {n}" for n in notes)
    return _text(slide, left, top, width, height, body, size=size, color=color, spacing=1.45)


def shot_slide(prs, eyebrow, title, image, notes, footer, page, total, *,
               caption=None):
    """A screenshot with a notes column beside it, house layout."""
    slide = _blank(prs)
    _head(slide, eyebrow, title)

    box_l, box_t, box_w = 0.51, 1.50, 8.24
    with Image.open(image) as im:
        ratio = im.height / im.width
    img_w = box_w - 0.08
    img_h = img_w * ratio
    max_h = 4.95
    if img_h > max_h:
        img_h = max_h
        img_w = img_h / ratio
    box_h = img_h + 0.08
    _rect(slide, box_l, box_t, img_w + 0.08, box_h, FRAME)
    slide.shapes.add_picture(image, Inches(box_l + 0.04), Inches(box_t + 0.04),
                             width=Inches(img_w), height=Inches(img_h))
    if caption:
        _text(slide, box_l + 0.04, box_t + box_h + 0.06, img_w, 0.3, caption,
              size=9.5, color=DIM, font=MONO)

    # The notes column is a fixed 3.78 x 4.60in, so long note sets have to shrink
    # to fit rather than run off the bottom of the slide into the page number.
    # Roughly 46 characters per line at 12.5pt in that width.
    lines = sum(max(1, -(-len(n) // 46)) for n in notes) + len(notes) * 0.35
    size = 12.5 if lines <= 15 else 11.5 if lines <= 18 else 10.5
    _bullets(slide, 9.02, 1.62, 3.78, 4.60, notes, size=size)
    _foot(slide, footer, page, total)
    return slide


def cards_slide(prs, eyebrow, title, cards, footer, page, total, *, note=None):
    """Numbered cards in a row -- the house 'concept' layout."""
    slide = _blank(prs)
    _head(slide, eyebrow, title)
    n = len(cards)
    gap, left0, total_w = 0.30, 0.55, 12.25
    w = (total_w - gap * (n - 1)) / n
    for i, (num, head, body) in enumerate(cards):
        x = left0 + i * (w + gap)
        _rect(slide, x, 1.75, w, 3.55, FRAME)
        _text(slide, x + 0.28, 1.98, w - 0.5, 0.4, num, size=17, color=AMBER,
              font=MONO, bold=True)
        _text(slide, x + 0.28, 2.48, w - 0.5, 0.5, head, size=15.5, color=TEXT, bold=True)
        _text(slide, x + 0.28, 3.05, w - 0.56, 2.0, body, size=11.5, color=SECOND,
              spacing=1.35)
    if note:
        _text(slide, 0.55, 5.60, 12.2, 0.9, note, size=11.5, color=AMBER, spacing=1.4)
    _foot(slide, footer, page, total)
    return slide


def table_slide(prs, eyebrow, title, headers, rows, footer, page, total, *,
                note=None, col_w=None, highlight=None):
    """A comparison table. `highlight` is a row index drawn in amber."""
    slide = _blank(prs)
    _head(slide, eyebrow, title)

    left, top, width = 0.55, 1.70, 12.25
    row_h = 0.42
    n_col = len(headers)
    widths = col_w or [width / n_col] * n_col
    scale = width / sum(widths)
    widths = [w * scale for w in widths]

    _rect(slide, left, top, width, row_h, FRAME, shape=MSO_SHAPE.RECTANGLE)
    x = left
    for w, h in zip(widths, headers):
        _text(slide, x + 0.14, top + 0.07, w - 0.2, 0.3, h, size=10.5, color=AMBER, font=MONO)
        x += w

    y = top + row_h + 0.06
    for r_i, row in enumerate(rows):
        is_hi = highlight is not None and r_i == highlight
        if is_hi:
            _rect(slide, left, y - 0.04, width, row_h, FRAME, shape=MSO_SHAPE.RECTANGLE)
        x = left
        for c_i, (w, cell) in enumerate(zip(widths, row)):
            _text(slide, x + 0.14, y, w - 0.2, 0.34, cell,
                  size=11.5,
                  color=AMBER if is_hi and c_i == 0 else (TEXT if c_i == 0 else SECOND),
                  font=SANS if c_i == 0 else MONO,
                  bold=is_hi and c_i == 0)
            x += w
        y += row_h

    if note:
        _text(slide, 0.55, min(y + 0.25, 6.10), 12.2, 1.0, note, size=11.5, color=AMBER,
              spacing=1.4)
    _foot(slide, footer, page, total)
    return slide


def prose_slide(prs, eyebrow, title, blocks, footer, page, total):
    """Two columns of headed paragraphs -- for the 'what broke and why' slides."""
    slide = _blank(prs)
    _head(slide, eyebrow, title)
    col_w, gap = 5.95, 0.35
    rows = (len(blocks) + 1) // 2
    card_h = min(2.30, (4.75 - 0.25 * (rows - 1)) / rows)
    for i, (head, body) in enumerate(blocks):
        col = i % 2
        row = i // 2
        x = 0.55 + col * (col_w + gap)
        y = 1.68 + row * (card_h + 0.25)
        _rect(slide, x, y, col_w, card_h, FRAME)
        _text(slide, x + 0.26, y + 0.17, col_w - 0.5, 0.35, head, size=13, color=AMBER, bold=True)
        _text(slide, x + 0.26, y + 0.62, col_w - 0.52, card_h - 0.75, body, size=11.5,
              color=SECOND, spacing=1.35)
    _foot(slide, footer, page, total)
    return slide


def two_shot_slide(prs, eyebrow, title, left_img, left_cap, right_img, right_cap,
                   notes, footer, page, total):
    """Two screenshots side by side, for before/after and OFF/ON."""
    slide = _blank(prs)
    _head(slide, eyebrow, title)
    box_w = 5.75
    for i, (img, cap) in enumerate(((left_img, left_cap), (right_img, right_cap))):
        x = 0.55 + i * (box_w + 0.30)
        with Image.open(img) as im:
            ratio = im.height / im.width
        img_w = box_w - 0.08
        img_h = img_w * ratio
        _rect(slide, x, 1.62, img_w + 0.08, img_h + 0.08, FRAME)
        slide.shapes.add_picture(img, Inches(x + 0.04), Inches(1.66),
                                 width=Inches(img_w), height=Inches(img_h))
        _text(slide, x + 0.04, 1.62 + img_h + 0.16, img_w, 0.3, cap, size=10,
              color=AMBER, font=MONO)
    _bullets(slide, 0.55, 5.05, 12.25, 1.75, notes, size=11.5)
    _foot(slide, footer, page, total)
    return slide


def closing_slide(prs, eyebrow, title, lines, footer, page, total):
    slide = _blank(prs)
    _head(slide, eyebrow, title)
    _rect(slide, 0.55, 1.75, 12.25, 4.45, FRAME)
    _text(slide, 0.95, 2.05, 11.5, 3.9,
          "\n".join(f"▪  {line}" for line in lines),
          size=13, color=SECOND, spacing=1.75)
    _foot(slide, footer, page, total)
    return slide
